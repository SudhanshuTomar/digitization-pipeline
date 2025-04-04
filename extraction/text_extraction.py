import os
import shutil
import zipfile
import tempfile
from io import BytesIO
import pandas as pd
import numpy as np

import dataiku
from docx import Document
import pdfplumber
from pptx import Presentation
import openpyxl
from PIL import Image
import pytesseract


class BaseExtractor:
    """
    Base extractor that handles file reading from a Dataiku Folder.
    """
    def __init__(self, folder_id):
        self.data_source = dataiku.Folder(folder_id) if folder_id else None

    def get_file_data(self, file_path):
        """
        Reads file data from the Dataiku Folder.
        """
        if self.data_source:
            with self.data_source.get_download_stream(file_path) as f:
                return f.read()
        return None
            
    def extract_file_to_temp(self, file_path):
        """
        Extracts a file to a temporary location and returns the temp path.
        """
        file_data = self.get_file_data(file_path)
        file_name = os.path.basename(file_path)
        temp_dir = tempfile.mkdtemp()
        temp_file_path = os.path.join(temp_dir, file_name)
        
        with open(temp_file_path, "wb") as temp_file:
            temp_file.write(file_data)
            
        return temp_file_path, temp_dir


class TextFileExtractor(BaseExtractor):
    """
    Extracts text from plain text files.
    """
    def extract_data(self, file_path, output_folder_id=None):
        """
        Extracts text from a plain text file.
        """
        try:
            file_data = self.get_file_data(file_path)
            # Decode the bytes to string with error handling
            text_content = file_data.decode('utf-8', errors='replace')
            
            return {
                "text_data": text_content,
                "table_data": [],
                "image_data": "",
                "images_extracted": []
            }
        except Exception as e:
            return {
                "text_data": f"Error extracting data from text file: {e}",
                "table_data": [],
                "image_data": "",
                "images_extracted": []
            }


class ImageProcessor:
    """
    Class for handling image processing, including OCR.
    """
    def __init__(self, lang="eng"):
        self.lang = lang
        
    def extract_text_from_image(self, image_data):
        """
        Extracts text from an image using OCR.
        """
        try:
            img = Image.open(BytesIO(image_data))
            text = pytesseract.image_to_string(img, lang=self.lang)
            return text.strip()
        except Exception as e:
            return f"Error extracting text from image: {e}"

    def extract_text_from_image_path(self, image_path):
        """
        Extracts text from an image file using OCR.
        """
        try:
            img = Image.open(image_path)
            text = pytesseract.image_to_string(img, lang=self.lang)
            return text.strip()
        except Exception as e:
            return f"Error extracting text from image at {image_path}: {e}"


class ImageExtractor(BaseExtractor):
    """
    Extracts and processes images.
    """
    def __init__(self, folder_id, lang="eng"):
        super().__init__(folder_id)
        self.image_processor = ImageProcessor(lang=lang)
    
    def extract_text(self, file_path):
        """
        Extracts text from an image using OCR.
        """
        try:
            file_data = self.get_file_data(file_path)
            text = self.image_processor.extract_text_from_image(file_data)
            return text
        except Exception as e:
            return f"Error processing image: {e}"
    
    def save_image(self, image_data, output_folder_id, image_path):
        """
        Saves an image to the output folder.
        """
        try:
            output_folder = dataiku.Folder(output_folder_id)
            with output_folder.get_writer(image_path) as writer:
                writer.write(image_data)
            return True
        except Exception as e:
            print(f"Error saving image {image_path}: {e}")
            return False


class WordExtractor(BaseExtractor):
    """
    Extracts text, tables, and images from Word documents.
    """
    def __init__(self, folder_id, lang="eng"):
        super().__init__(folder_id)
        self.image_processor = ImageProcessor(lang=lang)
    
    def extract_data(self, file_path, output_folder_id=None):
        """
        Extracts all data from a Word document.
        """
        try:
            file_data = self.get_file_data(file_path)
            doc_stream = BytesIO(file_data)
            doc = Document(doc_stream)
            
            # Extract text
            text_content = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
            text_data = "\n".join(text_content)
            
            # Extract tables
            table_data = []
            for table in doc.tables:
                table_content = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_content.append(row_data)
                table_data.append(table_content)
            
            # Extract images if output folder is provided
            images_extracted = []
            image_data = []
            
            if output_folder_id:
                # Fixed: Use the correct constructor for ImageExtractor
                temp_image_extractor = ImageExtractor(folder_id=self.data_source.get_id(), lang=self.image_processor.lang)
                images_extracted = self._extract_images(file_path, output_folder_id)
                
                if isinstance(images_extracted, list) and images_extracted:
                    # Get OCR text from each image
                    output_folder = dataiku.Folder(output_folder_id)
                    for img_path in images_extracted:
                        try:
                            with output_folder.get_download_stream(img_path) as f:
                                img_data = f.read()
                            ocr_text = self.image_processor.extract_text_from_image(img_data)
                            image_data.append(f"{img_path}: {ocr_text}")
                        except Exception as e:
                            image_data.append(f"{img_path}: Error - {str(e)}")
            
            return {
                "text_data": text_data,
                "table_data": table_data,
                "image_data": "\n\n".join(image_data) if image_data else "",
                "images_extracted": images_extracted
            }
        except Exception as e:
            import traceback
            return {
                "text_data": f"Error extracting data from Word document: {e}\n{traceback.format_exc()}",
                "table_data": [],
                "image_data": "",
                "images_extracted": []
            }

    def _extract_images(self, file_path, output_folder_id):
        """
        Extracts images from a Word document and saves them to a managed folder.
        """
        try:
            file_name = os.path.basename(file_path)
            base_name = os.path.splitext(file_name)[0]
            
            # Get temp file path
            temp_file_path, temp_dir = self.extract_file_to_temp(file_path)
            
            # Get output folder
            output_folder = dataiku.Folder(output_folder_id)
            existing_images = set(output_folder.list_paths_in_partition())
            
            # Extract images
            images_extracted = []
            with zipfile.ZipFile(temp_file_path, "r") as docx_zip:
                for file_info in docx_zip.infolist():
                    if file_info.filename.startswith("word/media/"):
                        image_name = os.path.basename(file_info.filename)
                        image_path = f"{base_name}_{image_name}"
                        
                        # Skip if image already exists
                        if image_path in existing_images:
                            images_extracted.append(image_path)
                            continue
                            
                        # Save image
                        image_data = docx_zip.read(file_info.filename)
                        with output_folder.get_writer(image_path) as writer:
                            writer.write(image_data)
                        images_extracted.append(image_path)
            
            # Clean up temp directory
            shutil.rmtree(temp_dir, ignore_errors=True)
            return images_extracted
        except Exception as e:
            return f"Error extracting images: {e}"


class PDFExtractor(BaseExtractor):
    """
    Extracts text and tables from PDF documents.
    """
    def __init__(self, folder_id, lang="eng"):
        super().__init__(folder_id)
        self.image_processor = ImageProcessor(lang=lang)
    
    def extract_data(self, file_path, output_folder_id=None):
        """
        Extracts all data from a PDF file.
        """
        try:
            file_data = self.get_file_data(file_path)
            pdf_stream = BytesIO(file_data)
            
            # Extract text
            text_content = []
            table_data = []
            images_extracted = []
            image_text = []
            
            with pdfplumber.open(pdf_stream) as pdf:
                for i, page in enumerate(pdf.pages):
                    # Extract text
                    extracted_text = page.extract_text()
                    if extracted_text:
                        text_content.append(extracted_text)
                    
                    # Extract tables
                    tables = page.extract_tables()
                    if tables:
                        table_data.extend(tables)
                    
                    # Extract images if output folder is provided
                    if output_folder_id:
                        page_images = page.images
                        if page_images:
                            for j, img in enumerate(page_images):
                                try:
                                    # Create image name
                                    file_name = os.path.basename(file_path)
                                    base_name = os.path.splitext(file_name)[0]
                                    img_path = f"{base_name}_page{i+1}_img{j+1}.png"
                                    
                                    # Extract and save image
                                    output_folder = dataiku.Folder(output_folder_id)
                                    
                                    # Skip if image already exists
                                    if img_path in output_folder.list_paths_in_partition():
                                        images_extracted.append(img_path)
                                        continue
                                    
                                    # Get image and save
                                    img_obj = page.to_image()
                                    img_data = img_obj.original
                                    
                                    # Save image
                                    with output_folder.get_writer(img_path) as writer:
                                        writer.write(img_data)
                                    
                                    # Extract text from image
                                    ocr_text = self.image_processor.extract_text_from_image(img_data)
                                    image_text.append(f"{img_path}: {ocr_text}")
                                    images_extracted.append(img_path)
                                except Exception as e:
                                    print(f"Error extracting image from PDF: {e}")
            
            return {
                "text_data": "\n\n".join(text_content),
                "table_data": table_data,
                "image_data": "\n\n".join(image_text) if image_text else "",
                "images_extracted": images_extracted
            }
        except Exception as e:
            return {
                "text_data": f"Error extracting data from PDF: {e}",
                "table_data": [],
                "image_data": "",
                "images_extracted": []
            }


class PowerPointExtractor(BaseExtractor):
    """
    Extracts text and images from PowerPoint presentations.
    """
    def __init__(self, folder_id, lang="eng"):
        super().__init__(folder_id)
        self.image_processor = ImageProcessor(lang=lang)
    
    def extract_data(self, file_path, output_folder_id=None):
        """
        Extracts all data from a PowerPoint file.
        """
        try:
            file_data = self.get_file_data(file_path)
            ppt_stream = BytesIO(file_data)
            prs = Presentation(ppt_stream)
            
            # Extract text
            text_content = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text.strip())
                
                text_content.append(f"Slide {i+1}: {' '.join(slide_text)}")
            
            # Extract images if output folder is provided
            images_extracted = []
            image_data = []
            
            if output_folder_id:
                # Extract images from the presentation
                file_name = os.path.basename(file_path)
                base_name = os.path.splitext(file_name)[0]
                temp_file_path, temp_dir = self.extract_file_to_temp(file_path)
                
                # Get output folder
                output_folder = dataiku.Folder(output_folder_id)
                existing_images = set(output_folder.list_paths_in_partition())
                
                try:
                    with zipfile.ZipFile(temp_file_path, "r") as pptx_zip:
                        for i, file_info in enumerate(pptx_zip.infolist()):
                            if file_info.filename.startswith("ppt/media/"):
                                image_name = os.path.basename(file_info.filename)
                                image_path = f"{base_name}_{image_name}"
                                
                                # Skip if image already exists
                                if image_path in existing_images:
                                    images_extracted.append(image_path)
                                    continue
                                
                                # Save image
                                image_data_bytes = pptx_zip.read(file_info.filename)
                                with output_folder.get_writer(image_path) as writer:
                                    writer.write(image_data_bytes)
                                images_extracted.append(image_path)
                                
                                # Extract text from image
                                try:
                                    ocr_text = self.image_processor.extract_text_from_image(image_data_bytes)
                                    image_data.append(f"{image_path}: {ocr_text}")
                                except Exception as e:
                                    image_data.append(f"{image_path}: Error - {str(e)}")
                except Exception as e:
                    print(f"Error extracting images from PPTX: {e}")
                
                # Clean up temp directory
                shutil.rmtree(temp_dir, ignore_errors=True)
            
            return {
                "text_data": "\n\n".join(text_content),
                "table_data": [],  # PowerPoint tables not extracted directly
                "image_data": "\n\n".join(image_data) if image_data else "",
                "images_extracted": images_extracted
            }
        except Exception as e:
            return {
                "text_data": f"Error extracting data from PowerPoint: {e}",
                "table_data": [],
                "image_data": "",
                "images_extracted": []
            }


class ExcelExtractor(BaseExtractor):
    """
    Extracts data from Excel files.
    """
    def extract_data(self, file_path, output_folder_id=None):
        """
        Extracts data from an Excel file.
        """
        try:
            file_data = self.get_file_data(file_path)
            excel_stream = BytesIO(file_data)
            
            # Use pandas to read Excel file
            all_sheets = pd.read_excel(excel_stream, sheet_name=None)
            
            # Extract text and table data
            text_content = []
            table_data = []
            
            for sheet_name, df in all_sheets.items():
                # Add sheet name to text content
                text_content.append(f"Sheet: {sheet_name}")
                
                # Convert DataFrame to list for table data
                table_data.append({
                    "sheet_name": sheet_name,
                    "headers": df.columns.tolist(),
                    "data": df.values.tolist()
                })
                
                # Add non-numeric cell values to text content
                for col in df.columns:
                    if df[col].dtype == object:  # Only process text columns
                        for value in df[col].dropna().unique():
                            if isinstance(value, str) and len(value) > 3:  # Only add meaningful text
                                text_content.append(str(value))
            
            return {
                "text_data": "\n\n".join(text_content),
                "table_data": table_data,
                "image_data": "",  # Excel images not extracted
                "images_extracted": []  # Excel images not extracted
            }
        except Exception as e:
            return {
                "text_data": f"Error extracting data from Excel: {e}",
                "table_data": [],
                "image_data": "",
                "images_extracted": []
            }


class FileProcessor:
    """
    Main processor for handling different file types.
    """
    def __init__(self, folder_id, output_folder_id=None, lang="eng"):
        self.folder_id = folder_id
        self.output_folder_id = output_folder_id
        self.lang = lang
        
        # Initialize extractors
        self.word_extractor = WordExtractor(folder_id, lang)
        self.pdf_extractor = PDFExtractor(folder_id, lang)
        self.ppt_extractor = PowerPointExtractor(folder_id, lang)
        self.excel_extractor = ExcelExtractor(folder_id)
        self.image_extractor = ImageExtractor(folder_id, lang)
        self.text_extractor = TextFileExtractor(folder_id)  # Added text file extractor
    
    def process_file(self, file_path):
        """
        Process a file based on its extension.
        """
        file_name = os.path.basename(file_path)
        ext = os.path.splitext(file_name)[1].lower()
        
        # Initialize result with file name
        result = {"file_name": file_name}
        
        try:
            # Process based on file type
            if ext == ".docx" or ext == ".doc":
                data = self.word_extractor.extract_data(file_path, self.output_folder_id)
            elif ext == ".pdf":
                data = self.pdf_extractor.extract_data(file_path, self.output_folder_id)
            elif ext in [".ppt", ".pptx"]:
                data = self.ppt_extractor.extract_data(file_path, self.output_folder_id)
            elif ext in [".xls", ".xlsx", ".xlsm"]:
                data = self.excel_extractor.extract_data(file_path)
            elif ext in [".jpg", ".jpeg", ".png", ".tiff", ".bmp", ".gif"]:
                text = self.image_extractor.extract_text(file_path)
                data = {
                    "text_data": "",
                    "table_data": [],
                    "image_data": text,
                    "images_extracted": [file_path]
                }
            elif ext in [".txt", ".csv", ".md", ".json", ".xml", ".html"]:  # Added text file formats
                data = self.text_extractor.extract_data(file_path)
            else:
                data = {
                    "text_data": f"Unsupported file type: {file_name}",
                    "table_data": [],
                    "image_data": "",
                    "images_extracted": []
                }
            
            # Update result with extracted data
            result.update({
                "text_data": data["text_data"],
                "table_data": str(data["table_data"]),  # Convert to string for DataFrame storage
                "image_data": data["image_data"],
                "images_extracted": str(data["images_extracted"])
            })
            
            return result
        except Exception as e:
            import traceback
            return {
                "file_name": file_name,
                "text_data": f"Error processing file: {e}\n{traceback.format_exc()}",
                "table_data": "[]",
                "image_data": "",
                "images_extracted": "[]"
            }
    
    def process_all_files(self, file_list):
        """
        Process all files in the list and return a DataFrame.
        """
        extracted_data = []
        for file_path in file_list:
            data = self.process_file(file_path)
            extracted_data.append(data)
        return pd.DataFrame(extracted_data)

